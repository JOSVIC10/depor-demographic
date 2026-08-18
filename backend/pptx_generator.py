import os
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from backend.models import Player, Team, Match, calculate_age
from backend.layout_engine import LayoutEngine, PitchSpec, calculate_box_width_and_font

# Colors
NAVY = RGBColor(0, 32, 96)         # #002060
PALE_GREEN = RGBColor(217, 234, 211) # #D9EAD3
DARK_GREEN = RGBColor(56, 118, 29)   # #38761D
PEACH = RGBColor(252, 229, 205)    # #FCE5CD
GOLD = RGBColor(212, 175, 55)      # #D4AF37
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(80, 80, 80)
RED_ARROW = RGBColor(204, 0, 0)
GREEN_ARROW = RGBColor(0, 153, 76)

def add_header_and_footer(slide, title_text: str, kpi_text: Optional[str], team: Team):
    # Header Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = "Calibri"

    # Header KPI Box (Dotted border)
    if kpi_text:
        kpi_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(0.35), Inches(4.0), Inches(0.7))
        kpi_box.fill.background()
        kpi_box.line.color.rgb = NAVY
        kpi_box.line.width = Pt(1.5)
        # Note: python-pptx doesn't have native dash line enum, but solid navy looks crisp
        tf_kpi = kpi_box.text_frame
        tf_kpi.word_wrap = True
        p_kpi = tf_kpi.paragraphs[0]
        p_kpi.text = kpi_text
        p_kpi.alignment = PP_ALIGN.CENTER
        p_kpi.font.size = Pt(12)
        p_kpi.font.bold = True
        p_kpi.font.color.rgb = NAVY

    # Footer Divider Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Footer Text
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.4))
    tf_f = footer_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = f"{team.club_name}  |  {team.season}  |  MEDICAL & SPORTS SCIENCE DEPARTMENT"
    p_f.alignment = PP_ALIGN.CENTER
    p_f.font.size = Pt(11)
    p_f.font.bold = True
    p_f.font.color.rgb = GOLD

def draw_pitch_background(slide, pitch_spec: PitchSpec):
    # Draw trapezoid field shape representing soccer pitch
    # Coordinates of trapezoid
    x_c = pitch_spec.x_center
    y_t = pitch_spec.y_top
    h = pitch_spec.height
    wt = pitch_spec.w_top
    wb = pitch_spec.w_bottom
    
    # We can create a field background box or load pitch image if available
    field = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(x_c - wb/2.0), Inches(y_t), Inches(wb), Inches(h))
    field.fill.solid()
    field.fill.fore_color.rgb = RGBColor(235, 245, 235)
    field.line.color.rgb = DARK_GREEN
    field.line.width = Pt(2.0)
    
    # Inner pitch line markings (halfway line, penalty boxes)
    line_mid = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x_c - pitch_spec.get_pitch_width_at_y(y_t + h/2.0)/2.0), Inches(y_t + h/2.0), Inches(pitch_spec.get_pitch_width_at_y(y_t + h/2.0)), Inches(0.02))
    line_mid.fill.solid()
    line_mid.fill.fore_color.rgb = DARK_GREEN
    line_mid.line.fill.background()

class PPTXGenerator:
    def __init__(self, template_path: Optional[str] = None):
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            # Clear existing slides if using reference deck as base layout
            for i in range(len(self.prs.slides)-1, -1, -1):
                rId = self.prs.slides._sldIdLst[i].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[i]
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)

        self.layout_engine = LayoutEngine()

    def generate_demographic_slide(self, team: Team, players: List[Player]):
        blank_slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        title = f"Analysis | Player\n{team.name} SQUAD DEMOGRAPHIC"
        kpi = f"TOTAL PLAYERS: {len(players)}\nAVG AGE: {round(sum(p.age for p in players)/len(players), 1) if players else 0.0}".replace('.', ',')
        add_header_and_footer(slide, title, kpi, team)

        # Categorize players by age band & position
        age_bands = [
            ("30+", lambda a: a >= 30),
            ("26-29", lambda a: 26 <= a <= 29),
            ("22-25", lambda a: 22 <= a <= 25),
            ("18-21", lambda a: 18 <= a <= 21),
        ]
        
        categories = ["Porteros", "Centrales", "Laterales", "Mediocentros", "Int/Extremos", "Delanteros"]
        
        # Grid: row -> cat -> list of players
        grid = {b_label: {cat: [] for cat in categories} for b_label, _ in age_bands}
        
        for p in players:
            p_cat = p.derived_category if p.derived_category in categories else "Mediocentros"
            for b_label, condition in age_bands:
                if condition(p.age):
                    grid[b_label][p_cat].append(p)
                    break

        # Calculate dynamic column widths
        col_widths = self.layout_engine.calculate_demographic_column_widths([p.model_dump() for p in players], total_table_width=12.2)
        
        # Add Table
        rows = 7 # Header + 4 bands + TOTALS + MEDIA EDAD
        cols = 8 # Age Band + 6 categories + TOTAL
        
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.35), Inches(12.2), Inches(5.3))
        table = table_shape.table
        
        # Set column widths
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = Inches(w)
            
        headers = ["BANDA DE EDAD"] + [c.upper() for c in categories] + ["TOTAL"]
        for idx, h in enumerate(headers):
            cell = table.cell(0, idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            p = cell.text_frame.paragraphs[0]
            p.text = h
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE

        # Populate Age Bands
        col_totals = [0] * len(categories)
        col_ages_sum = [0.0] * len(categories)

        for r_idx, (b_label, _) in enumerate(age_bands, start=1):
            cell_band = table.cell(r_idx, 0)
            cell_band.fill.solid()
            cell_band.fill.fore_color.rgb = LIGHT_GRAY
            p = cell_band.text_frame.paragraphs[0]
            p.text = b_label
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = NAVY
            
            row_total = 0
            for c_idx, cat in enumerate(categories, start=1):
                cell = table.cell(r_idx, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                p_list = grid[b_label][cat]
                row_total += len(p_list)
                col_totals[c_idx-1] += len(p_list)
                for p_item in p_list:
                    col_ages_sum[c_idx-1] += p_item.age

                tf = cell.text_frame
                tf.word_wrap = True
                tf.clear()
                if p_list:
                    for p_item in p_list:
                        p_txt = tf.add_paragraph()
                        p_txt.text = f"{p_item.name} - {p_item.age}"
                        p_txt.font.size = Pt(8.5)
                        p_txt.font.color.rgb = NAVY
                else:
                    p_txt = tf.paragraphs[0]
                    p_txt.text = "-"
                    p_txt.alignment = PP_ALIGN.CENTER
                    p_txt.font.size = Pt(9)
                    p_txt.font.color.rgb = DARK_GRAY

            # Row Total cell
            cell_tot = table.cell(r_idx, 7)
            cell_tot.fill.solid()
            cell_tot.fill.fore_color.rgb = LIGHT_GRAY
            p = cell_tot.text_frame.paragraphs[0]
            p.text = str(row_total)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = NAVY

        # Row 5: TOTALS
        cell_t_label = table.cell(5, 0)
        cell_t_label.fill.solid()
        cell_t_label.fill.fore_color.rgb = NAVY
        p = cell_t_label.text_frame.paragraphs[0]
        p.text = "TOTALS"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE

        grand_total = sum(col_totals)
        for c_idx in range(1, 7):
            cell = table.cell(5, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALE_GREEN
            p = cell.text_frame.paragraphs[0]
            p.text = str(col_totals[c_idx-1])
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = NAVY

        cell_gt = table.cell(5, 7)
        cell_gt.fill.solid()
        cell_gt.fill.fore_color.rgb = NAVY
        p = cell_gt.text_frame.paragraphs[0]
        p.text = str(grand_total)
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

        # Row 6: MEDIA EDAD
        cell_m_label = table.cell(6, 0)
        cell_m_label.fill.solid()
        cell_m_label.fill.fore_color.rgb = NAVY
        p = cell_m_label.text_frame.paragraphs[0]
        p.text = "MEDIA EDAD"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE

        for c_idx in range(1, 7):
            cell = table.cell(6, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
            count = col_totals[c_idx-1]
            avg = (col_ages_sum[c_idx-1] / count) if count > 0 else 0.0
            p = cell.text_frame.paragraphs[0]
            p.text = f"{avg:.1f}".replace('.', ',') if count > 0 else "-"
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(10)
            p.font.color.rgb = NAVY

        all_avg = (sum(col_ages_sum) / grand_total) if grand_total > 0 else 0.0
        cell_gavg = table.cell(6, 7)
        cell_gavg.fill.solid()
        cell_gavg.fill.fore_color.rgb = NAVY
        p = cell_gavg.text_frame.paragraphs[0]
        p.text = f"{all_avg:.1f}".replace('.', ',') if grand_total > 0 else "-"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE

    def generate_squad_pitch_slide(self, team: Team, players: List[Player]):
        blank_slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        title = f"Analysis | Player\n{team.name}"
        kpi = f"NUMERO DE JUGADORES: {len(players)} JUGADORES"
        add_header_and_footer(slide, title, kpi, team)

        pitch_spec = PitchSpec()
        draw_pitch_background(slide, pitch_spec)

        boxes = self.layout_engine.layout_full_squad([p.model_dump() for p in players])

        for box in boxes:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box.x), Inches(box.y), Inches(box.width), Inches(box.height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PALE_GREEN
            shape.line.color.rgb = DARK_GREEN
            shape.line.width = Pt(1.2)
            
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = box.label
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(box.font_size)
            p.font.size = Pt(box.font_size)
            p.font.color.rgb = NAVY

        # Render injured players (Left)
        injured_players = [p for p in players if getattr(p, 'is_injured', False)]
        if injured_players:
            txBox = slide.shapes.add_textbox(Inches(0.2), Inches(1.2), Inches(2.5), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "LESIONADOS / INJURED"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = NAVY
            
            for ip in injured_players:
                p = tf.add_paragraph()
                extra_parts = []
                if getattr(ip, 'injury_description', None):
                    extra_parts.append(str(ip.injury_description).strip())
                if getattr(ip, 'injury_phase', None):
                    extra_parts.append(str(ip.injury_phase).strip())
                if getattr(ip, 'injury_return_time', None):
                    extra_parts.append(str(ip.injury_return_time).strip())
                info_text = f" ({' | '.join(extra_parts)})" if extra_parts else ""
                p.text = f"• {ip.name}{info_text}"
                p.font.size = Pt(8.5)
                p.font.color.rgb = NAVY

        # Render extra/filial players (Right)
        extra_filial_players = [p for p in players if getattr(p, 'extra_pitch_team_id', None) == team.id and p.team_id != team.id]
        if extra_filial_players:
            txBoxRight = slide.shapes.add_textbox(Inches(10.2), Inches(1.2), Inches(2.8), Inches(0.5))
            tfRight = txBoxRight.text_frame
            p = tfRight.paragraphs[0]
            p.text = "FILIAL / INCORPORADOS"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = NAVY
            
            for ep in extra_filial_players:
                p = tfRight.add_paragraph()
                team_label = "FABRIL" if ep.team_id == "fabril" else ("JUVENIL A" if ep.team_id == "juvenil_a" else ep.team_id.upper())
                p.text = f"• {ep.name} [{team_label}] ({ep.detailed_position})"
                p.font.size = Pt(8.5)
                p.font.color.rgb = NAVY

    def generate_match_report_slide(self, full_match: Dict):
        match = full_match['match']
        team = full_match['team']
        starters = full_match['starters']
        substitutes = full_match['substitutes']
        subs_events = full_match['substitutions']
        players_map = full_match['players_map']

        blank_slide_layout = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        title = f"Analysis | Player\n{match.title_display}"
        kpi = f"TOTAL NUMBER OF SUBSTITUTIONS: {len(subs_events)}"
        add_header_and_footer(slide, title, kpi, team)

        # Pitch in center
        pitch_spec = PitchSpec(x_center=6.4, y_top=1.35, height=5.25, w_top=5.5, w_bottom=7.8)
        draw_pitch_background(slide, pitch_spec)

        # Map substitutions
        subbed_out_map = {se.player_out_id: se for se in subs_events}
        subbed_in_map = {se.player_in_id: se for se in subs_events}

        starters_dicts = []
        for s in starters:
            sd = s.model_dump()
            if s.player_id in subbed_out_map:
                sd['substituted_minute'] = subbed_out_map[s.player_id].minute
            starters_dicts.append(sd)

        layout_eng = LayoutEngine(pitch_spec)
        starter_boxes = layout_eng.layout_match_starters(starters_dicts, {k: v.model_dump() for k, v in players_map.items()})

        for box in starter_boxes:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box.x), Inches(box.y), Inches(box.width), Inches(box.height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PALE_GREEN
            shape.line.color.rgb = DARK_GREEN
            shape.line.width = Pt(1.5)
            
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = box.label
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p.font.size = Pt(box.font_size)
            p.font.color.rgb = NAVY

            # Draw Red Down Arrow if subbed off
            if box.id in subbed_out_map:
                arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(box.x + box.width - 0.15), Inches(box.y + 0.05), Inches(0.18), Inches(0.22))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RED_ARROW
                arrow.line.fill.background()

        # Left Panel: Substitutes List (Peach background)
        sub_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.35), Inches(2.3), Inches(5.25))
        sub_box.fill.solid()
        sub_box.fill.fore_color.rgb = PEACH
        sub_box.line.color.rgb = RGBColor(226, 170, 110)
        
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p = tf_sub.paragraphs[0]
        p.text = "Substitutes"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY

        for sub_entry in substitutes:
            p_obj = players_map.get(sub_entry.player_id)
            p_name = p_obj.name if p_obj else "Suplente"
            has_yellow = getattr(sub_entry, 'has_yellow_card', False)
            has_red = getattr(sub_entry, 'has_red_card', False)
            c_min = getattr(sub_entry, 'card_minute', None)
            goals = getattr(sub_entry, 'goals', 0)
            goal_badge = f" (⚽ {goals})" if goals and goals > 0 else ""
            card_badge = f" ({c_min}’ 🟥)" if (has_red and c_min) else (" 🟥" if has_red else (f" ({c_min}’ 🟨)" if (has_yellow and c_min) else (" 🟨" if has_yellow else "")))

            p_txt = tf_sub.add_paragraph()
            
            if sub_entry.player_id in subbed_in_map:
                ev = subbed_in_map[sub_entry.player_id]
                p_txt.text = f"▲ {p_name} ({ev.minute}’){goal_badge}{card_badge}"
                p_txt.font.bold = True
                p_txt.font.size = Pt(9.5)
                p_txt.font.color.rgb = GREEN_ARROW
            else:
                p_txt.text = f"• {p_name}{goal_badge}{card_badge}"
                p_txt.font.size = Pt(9.0)
                p_txt.font.color.rgb = NAVY

        # Right Panel: Substitutions Log & Cadence
        log_box = slide.shapes.add_textbox(Inches(10.6), Inches(1.35), Inches(2.3), Inches(5.25))
        tf_log = log_box.text_frame
        tf_log.word_wrap = True
        
        p = tf_log.paragraphs[0]
        p.text = "Substitutions"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY

        for idx, ev in enumerate(subs_events, start=1):
            p_out = players_map.get(ev.player_out_id)
            p_in = players_map.get(ev.player_in_id)
            out_name = p_out.name if p_out else "Out"
            in_name = p_in.name if p_in else "In"
            
            in_sub = next((s for s in substitutes if s.player_id == ev.player_in_id), None)
            in_goals = getattr(in_sub, 'goals', 0) if in_sub else 0
            in_goal_str = f" ⚽{in_goals}" if in_goals > 0 else ""
            in_card_str = " 🟥" if (in_sub and getattr(in_sub, 'has_red_card', False)) else (" 🟨" if (in_sub and getattr(in_sub, 'has_yellow_card', False)) else "")
            
            p_item = tf_log.add_paragraph()
            p_item.text = f"{idx}. Min {ev.minute}’:\n   ▲ {in_name}{in_goal_str}{in_card_str}\n   ▼ (sale {out_name})"
            p_item.font.size = Pt(9.0)
            p_item.font.color.rgb = NAVY

    def save(self, output_path: str):
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        self.prs.save(output_path)
